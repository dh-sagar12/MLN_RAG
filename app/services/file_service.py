"""File service for handling file uploads and text extraction."""
import os
import aiofiles
from pathlib import Path
from typing import Optional, Tuple
from starlette.datastructures import UploadFile
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import markdown
from html.parser import HTMLParser
from app.config import settings


class HTMLTextExtractor(HTMLParser):
    """Simple HTML text extractor."""
    
    def __init__(self):
        super().__init__()
        self.text = []
    
    def handle_data(self, data: str) -> None:
        self.text.append(data)
    
    def get_text(self) -> str:
        return " ".join(self.text)


class FileService:
    """Service for handling file operations."""
    
    def __init__(self):
        self.storage_path = Path(settings.storage_path)
        self.storage_path.mkdir(parents=True, exist_ok=True)
    
    async def save_file(
        self, 
        file: UploadFile, 
        kb_id: str
    ) -> Tuple[str, int]:
        """Save uploaded file to storage.
        
        Args:
            file: Uploaded file
            kb_id: Knowledge base ID
            
        Returns:
            Tuple of (file_path, file_size)
        """
        kb_path = self.storage_path / str(kb_id)
        kb_path.mkdir(parents=True, exist_ok=True)
        
        file_path = kb_path / file.filename
        file_size = 0
        
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)
            file_size = len(content)
        
        return str(file_path), file_size
    
    async def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from various file types.
        
        Args:
            file_path: Path to the file
            file_type: MIME type or file extension
            
        Returns:
            Extracted text content
        """
        path = Path(file_path)
        
        if file_type == "application/pdf" or path.suffix.lower() == ".pdf":
            return await self._extract_pdf(file_path)
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or path.suffix.lower() == ".docx":
            return await self._extract_docx(file_path)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or path.suffix.lower() == ".pptx":
            return await self._extract_pptx(file_path)
        elif file_type == "text/markdown" or path.suffix.lower() == ".md":
            return await self._extract_markdown(file_path)
        elif file_type == "text/html" or path.suffix.lower() == ".html":
            return await self._extract_html(file_path)
        elif file_type == "text/plain" or path.suffix.lower() == ".txt":
            return await self._extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    async def _extract_pdf(self, file_path: str) -> str:
        """
        Extract clean text from PDF using PyMuPDF.
        This avoids the broken word-by-word line issues seen in pypdf.
        """
        import fitz
        text_parts = []
        doc = fitz.open(file_path)

        for page in doc:
            # "text" mode preserves proper reading order & paragraphs
            text = page.get_text("text")
            if text:
                text_parts.append(text)

        full_text = "\n".join(text_parts)
        return self._normalize_text(full_text)

    def _normalize_text(self, text: str) -> str:
        """
        Optional light cleanup for PDFs with irregular spacing.
        Removes excessive empty lines, trims spaces, formats well.
        """
        lines = [line.strip() for line in text.splitlines()]
        cleaned = []

        for line in lines:
            # remove empty lines but keep paragraph breaks
            if line:
                cleaned.append(line)

        return "\n".join(cleaned)

    
    async def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        doc = DocxDocument(file_path)
        text_parts = [para.text for para in doc.paragraphs if para.text]

        # Include table contents (rows separated by newlines, cells by tabs)
        for table in doc.tables:
            for row in table.rows:
                cell_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cell_text:
                    text_parts.append("\t".join(cell_text))

        return "\n".join(text_parts)
    
    async def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    
    async def _extract_markdown(self, file_path: str) -> str:
        """Extract text from Markdown."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
            # Convert markdown to HTML then extract text
            html = markdown.markdown(content)
            parser = HTMLTextExtractor()
            parser.feed(html)
            return parser.get_text()
    
    async def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
            parser = HTMLTextExtractor()
            parser.feed(content)
            return parser.get_text()
    
    async def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            return await f.read()
    
    def get_file_type(self, filename: str, content_type: Optional[str] = None) -> str:
        """Get file type from filename or content type."""
        if content_type:
            return content_type
        
        ext = Path(filename).suffix.lower()
        type_map = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".md": "text/markdown",
            ".html": "text/html",
            ".txt": "text/plain",
        }
        return type_map.get(ext, "application/octet-stream")

